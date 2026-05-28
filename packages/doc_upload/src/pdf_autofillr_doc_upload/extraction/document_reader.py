# pdf_autofillr_doc_upload/extraction/document_reader.py
"""
DocumentReader — extract plain text from any supported document format.

Supported formats:
    .pdf    PyMuPDF
    .docx   python-docx
    .pptx   python-pptx
    .xlsx   openpyxl
    .xls    openpyxl
    .csv    built-in csv
    .json   built-in json
    .txt    plain read
    .md     plain read (Markdown)
    .html   html.parser strip tags
    .htm    html.parser strip tags
    .xml    ElementTree text walk
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

# ── Per-format extractors ──────────────────────────────────────────────────────


def _read_pdf(path: str) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages = []
    for page in doc:
        txt = page.get_text("text")
        if txt:
            pages.append(txt.strip())
    return "\n\n".join(pages)


def _read_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _read_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        content = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())
        if len(content) > 1:
            slides.append("\n".join(content))
    return "\n\n".join(slides)


def _read_xlsx(path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        lines = [f"--- Sheet: {name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row).strip()
            if row_text:
                lines.append(row_text)
        if len(lines) > 1:
            sheets.append("\n".join(lines))
    return "\n\n".join(sheets)


def _read_csv(path: str) -> str:
    lines = []
    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = " | ".join(cell.strip() for cell in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _read_json(path: str) -> str:
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_text(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read().strip()


def _read_html(path: str) -> str:
    from html.parser import HTMLParser

    class _Stripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts = []
            self._skip_tags = {"script", "style", "head"}
            self._in_skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in self._skip_tags:
                self._in_skip += 1

        def handle_endtag(self, tag):
            if tag in self._skip_tags and self._in_skip:
                self._in_skip -= 1

        def handle_data(self, data):
            if not self._in_skip and data.strip():
                self._parts.append(data.strip())

        def get_text(self):
            return "\n".join(self._parts)

    with open(path, encoding="utf-8", errors="replace") as f:
        content = f.read()
    s = _Stripper()
    s.feed(content)
    return s.get_text()


def _read_xml(path: str) -> str:
    import xml.etree.ElementTree as ET

    tree = ET.parse(path)
    parts = []

    def _walk(el):
        if el.text and el.text.strip():
            parts.append(el.text.strip())
        for child in el:
            _walk(child)
        if el.tail and el.tail.strip():
            parts.append(el.tail.strip())

    _walk(tree.getroot())
    return "\n".join(parts)


# ── Dispatch table ─────────────────────────────────────────────────────────────

_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
    ".xls": _read_xlsx,
    ".csv": _read_csv,
    ".json": _read_json,
    ".txt": _read_text,
    ".md": _read_text,
    ".markdown": _read_text,
    ".html": _read_html,
    ".htm": _read_html,
    ".xml": _read_xml,
}

SUPPORTED_EXTENSIONS = list(_READERS.keys())


class DocumentReader:
    """
    Detect file format and extract plain text.

    Usage::

        reader = DocumentReader()
        text = reader.read("/path/to/investor_profile.pdf")
    """

    def read(self, file_path: str) -> str:
        """
        Extract text from any supported file format.

        Args:
            file_path: Path to the document.

        Returns:
            Extracted plain text.

        Raises:
            FileNotFoundError: File does not exist.
            ValueError:        Unsupported file extension.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")

        ext = path.suffix.lower()
        reader_fn = _READERS.get(ext)
        if reader_fn is None:
            raise ValueError(
                f"Unsupported file format: {ext!r}. "
                f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
            )

        return reader_fn(str(path))

    @staticmethod
    def supported_extensions() -> list:
        return SUPPORTED_EXTENSIONS
